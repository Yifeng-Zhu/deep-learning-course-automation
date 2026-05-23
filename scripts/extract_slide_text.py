"""Extract visible slide text from original lecture PowerPoint decks.

The script reads the course configuration files, discovers PowerPoint files
using the "{lecture_number} - {lecture_title}.pptx" naming convention, and
writes a CSV slide-text inventory inside this repository. Batch mode processes
lecture numbers 1 and above; lecture 0 is a special case and is processed only
when explicitly requested.
"""

from __future__ import annotations

import argparse
import csv
import sys
from pathlib import Path
from typing import Any

from lecture_files import LectureDeck, discover_lecture_decks
from pptx_reader import PptxReadError, open_presentation_readonly

try:
    import yaml
    from yaml import YAMLError
except ImportError:  # pragma: no cover - depends on the local Python environment
    yaml = None

    class YAMLError(Exception):
        """Placeholder when PyYAML is not installed."""

try:
    from pptx.enum.shapes import PP_PLACEHOLDER
except ImportError:  # pragma: no cover - depends on the local Python environment
    PP_PLACEHOLDER = None


OUTPUT_CSV = Path("manifests") / "slide_text_inventory.csv"
CSV_FIELDS = [
    "lecture_number",
    "lecture_title",
    "source_filename",
    "slide_number",
    "possible_slide_title",
    "visible_text",
]
ERROR_CSV_FIELDS = [
    "lecture_number",
    "lecture_title",
    "source_filename",
    "error",
]


def load_yaml(path: Path) -> dict[str, Any]:
    if yaml is None:
        raise RuntimeError(
            "PyYAML is required. Install dependencies with: pip install -r requirements.txt"
        )

    with path.open("r", encoding="utf-8") as handle:
        data = yaml.safe_load(handle)

    if not isinstance(data, dict):
        raise ValueError(f"{path.name} must contain a YAML mapping.")

    return data


def get_project_root() -> Path:
    return Path(__file__).resolve().parents[1]


def get_original_slides_folder(project_root: Path) -> Path:
    shared_config = load_yaml(project_root / "course_config.yaml")
    local_config = load_yaml(project_root / "course_config.local.yaml")

    course_drive_root_value = local_config.get("course_drive_root")
    folders = shared_config.get("folders")

    if not isinstance(course_drive_root_value, str) or not course_drive_root_value.strip():
        raise ValueError("course_config.local.yaml is missing course_drive_root.")

    if not isinstance(folders, dict):
        raise ValueError("course_config.yaml is missing the folders mapping.")

    original_slides_name = folders.get("original_slides")
    if not isinstance(original_slides_name, str) or not original_slides_name.strip():
        raise ValueError("course_config.yaml is missing folders.original_slides.")

    return Path(course_drive_root_value).expanduser() / original_slides_name


def select_decks(
    decks: list[LectureDeck],
    lecture_number: int | None,
    include_lecture_zero: bool,
) -> list[LectureDeck]:
    if lecture_number is not None:
        matching_decks = [deck for deck in decks if deck.lecture_number == lecture_number]
        if not matching_decks:
            available = ", ".join(str(deck.lecture_number) for deck in decks) or "none"
            raise ValueError(
                f"Lecture {lecture_number} was not found. Available lecture numbers: {available}"
            )
        return matching_decks

    selected_decks = [
        deck for deck in decks if include_lecture_zero or deck.lecture_number >= 1
    ]
    if not selected_decks:
        raise ValueError("No batch-eligible lecture decks found. Expected lecture numbers 1 and above.")

    return selected_decks


def clean_text(value: str) -> str:
    lines = [" ".join(line.split()) for line in value.splitlines()]
    return "\n".join(line for line in lines if line)


def shape_xml_text(shape: Any) -> list[str]:
    text_parts: list[str] = []
    for node in shape.element.iter():
        if node.tag.endswith("}t") and node.text:
            cleaned = clean_text(node.text)
            if cleaned:
                text_parts.append(cleaned)
    return text_parts


def slide_visible_text(slide: Any) -> str:
    text_parts: list[str] = []
    for shape in slide.shapes:
        text_parts.extend(shape_xml_text(shape))
    return "\n".join(text_parts)


def shape_text(shape: Any) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return clean_text(shape.text_frame.text)


def possible_slide_title(slide: Any, visible_text: str) -> str:
    if PP_PLACEHOLDER is not None:
        title_types = {
            PP_PLACEHOLDER.TITLE,
            PP_PLACEHOLDER.CENTER_TITLE,
            PP_PLACEHOLDER.VERTICAL_TITLE,
        }
        for shape in slide.shapes:
            if not getattr(shape, "is_placeholder", False):
                continue
            if shape.placeholder_format.type in title_types:
                title = shape_text(shape)
                if title:
                    return title

    for line in visible_text.splitlines():
        candidate = line.strip()
        if candidate:
            return candidate

    return ""


def extract_slide_rows(deck: LectureDeck) -> list[dict[str, str | int]]:
    try:
        with open_presentation_readonly(deck.path) as (presentation, _repair_info):
            rows: list[dict[str, str | int]] = []

            for slide_index, slide in enumerate(presentation.slides, start=1):
                visible_text = slide_visible_text(slide)
                rows.append(
                    {
                        "lecture_number": deck.lecture_number,
                        "lecture_title": deck.lecture_title,
                        "source_filename": deck.path.name,
                        "slide_number": slide_index,
                        "possible_slide_title": possible_slide_title(slide, visible_text),
                        "visible_text": visible_text,
                    }
                )

            return rows
    except PptxReadError as exc:
        raise RuntimeError(
            f"Could not read Lecture {deck.lecture_number} PowerPoint file "
            f"{deck.path.name}: {exc}"
        ) from exc


def write_inventory(rows: list[dict[str, str | int]], output_path: Path) -> None:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    with output_path.open("w", encoding="utf-8-sig", newline="") as handle:
        writer = csv.DictWriter(handle, fieldnames=CSV_FIELDS)
        writer.writeheader()
        writer.writerows(rows)


def write_error_report(rows: list[dict[str, str | int]], output_path: Path) -> None:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    with output_path.open("w", encoding="utf-8-sig", newline="") as handle:
        writer = csv.DictWriter(handle, fieldnames=ERROR_CSV_FIELDS)
        writer.writeheader()
        writer.writerows(rows)


def error_report_path(output_path: Path) -> Path:
    return output_path.with_name(f"{output_path.stem}_errors{output_path.suffix}")


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Extract visible text from original lecture PowerPoint decks."
    )
    parser.add_argument(
        "--lecture-number",
        type=int,
        help="Process one lecture number explicitly. Use 0 only when Lecture 0 is requested.",
    )
    parser.add_argument(
        "--include-lecture-zero",
        action="store_true",
        help="Include Lecture 0 in batch extraction. By default, batch mode starts at Lecture 1.",
    )
    parser.add_argument(
        "--output",
        default=str(OUTPUT_CSV),
        help=f"Output CSV path relative to the repository root. Defaults to {OUTPUT_CSV}.",
    )
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    project_root = get_project_root()
    output_path = project_root / args.output
    error_path = error_report_path(output_path)

    try:
        original_slides_folder = get_original_slides_folder(project_root)
        decks = discover_lecture_decks(original_slides_folder)
        selected_decks = select_decks(
            decks=decks,
            lecture_number=args.lecture_number,
            include_lecture_zero=args.include_lecture_zero,
        )
    except (OSError, RuntimeError, ValueError, YAMLError) as exc:
        print("Slide text extraction failed.")
        print(f"Error: {exc}")
        return 1

    rows: list[dict[str, str | int]] = []
    error_rows: list[dict[str, str | int]] = []
    slides_by_lecture: dict[int, int] = {}

    for deck in selected_decks:
        try:
            deck_rows = extract_slide_rows(deck)
        except RuntimeError as exc:
            error_rows.append(
                {
                    "lecture_number": deck.lecture_number,
                    "lecture_title": deck.lecture_title,
                    "source_filename": deck.path.name,
                    "error": str(exc),
                }
            )
            continue

        rows.extend(deck_rows)
        slides_by_lecture[deck.lecture_number] = len(deck_rows)

    try:
        write_inventory(rows, output_path)
        write_error_report(error_rows, error_path)
    except (OSError, RuntimeError, ValueError, YAMLError) as exc:
        print("Slide text extraction failed.")
        print(f"Error: {exc}")
        return 1

    successful_decks = [
        deck for deck in selected_decks if deck.lecture_number in slides_by_lecture
    ]
    attempted_numbers = ", ".join(str(deck.lecture_number) for deck in selected_decks)
    processed_numbers = ", ".join(str(deck.lecture_number) for deck in successful_decks) or "none"
    attempted_files = ", ".join(deck.path.name for deck in selected_decks)
    skipped_lecture_zero = (
        args.lecture_number is None
        and not args.include_lecture_zero
        and any(deck.lecture_number == 0 for deck in decks)
    )

    if error_rows:
        print("Slide text extraction completed with errors")
    else:
        print("Slide text extraction complete")
    print(f"Original slides folder: {original_slides_folder}")
    print(f"Discovered decks:       {len(decks)}")
    if skipped_lecture_zero:
        print("Lecture 0 special case: skipped in batch mode")
    print(f"Lectures attempted:     {attempted_numbers}")
    print(f"Lectures processed:     {processed_numbers}")
    print(f"Decks attempted:        {len(selected_decks)}")
    print(f"Decks processed:        {len(successful_decks)}")
    print(f"Decks failed:           {len(error_rows)}")
    print(f"Slides processed:       {len(rows)}")
    for lecture_number in sorted(slides_by_lecture):
        print(f"  - Lecture {lecture_number}: {slides_by_lecture[lecture_number]} slides")
    if error_rows:
        for error_row in error_rows:
            print(
                f"  - FAILED Lecture {error_row['lecture_number']}: "
                f"{error_row['source_filename']}"
            )
    print(f"Source files attempted: {attempted_files}")
    print(f"Output CSV:             {output_path}")
    print(f"Error CSV:              {error_path}")
    print("PowerPoint safety: original decks were read only; no PowerPoint files were modified.")
    return 1 if error_rows else 0


if __name__ == "__main__":
    sys.exit(main())
