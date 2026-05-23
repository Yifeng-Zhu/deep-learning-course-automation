"""Generate revised PowerPoint decks for approved video segments.

The script reads the approved rows in manifests/video_segments/*.csv, opens the
matching original lecture deck from the configured 01_Original_Slides folder,
and writes one segment deck per approved video to 02_Revised_Slides/{lecture}/.

Original PowerPoint files are opened read-only and are never overwritten.
"""

from __future__ import annotations

import argparse
import csv
import re
import sys
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from lecture_files import LectureDeck, discover_lecture_decks
from pptx_reader import PptxReadError, open_presentation_readonly

try:
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import nsdecls
    import yaml
    from yaml import YAMLError
except ImportError:  # pragma: no cover - depends on the local Python environment
    parse_xml = None
    nsdecls = None
    yaml = None

    class YAMLError(Exception):
        """Placeholder when PyYAML is not installed."""

try:
    from pptx.enum.shapes import PP_PLACEHOLDER
    from pptx.util import Inches, Pt
except ImportError:  # pragma: no cover - surfaced by main()
    PP_PLACEHOLDER = None
    Inches = None
    Pt = None


VIDEO_SEGMENTS_DIR = Path("manifests") / "video_segments"
TEACHING_MAPS_DIR = Path("manifests") / "teaching_maps"
APPROVED_STATUS = "approved_for_pptx"
OUTPUT_FIELD = "output_pptx"
SCRIPT_MARKER = "Generated speaker-script outline:"


@dataclass(frozen=True)
class Segment:
    source_csv: Path
    row_index: int
    lecture_number: int
    lecture_title: str
    video_number: int
    video_title: str
    slide_start: int
    slide_end: int
    estimated_minutes: str
    main_concept_group: str
    learning_objectives: str
    logical_segment_reason: str
    stopping_point_reason: str
    recap: str
    quiz: str
    review_status: str


def project_root() -> Path:
    return Path(__file__).resolve().parents[1]


def load_yaml(path: Path) -> dict[str, Any]:
    if yaml is None:
        raise RuntimeError(
            "PyYAML is required. Install dependencies with: pip install -r requirements.txt"
        )

    with path.open("r", encoding="utf-8") as handle:
        data = yaml.safe_load(handle)

    if not isinstance(data, dict):
        raise ValueError(f"{path.name} must contain a YAML mapping.")

    return data


def configured_folder(project_root_path: Path, folder_key: str) -> Path:
    shared_config = load_yaml(project_root_path / "course_config.yaml")
    local_config = load_yaml(project_root_path / "course_config.local.yaml")

    course_drive_root_value = local_config.get("course_drive_root")
    folders = shared_config.get("folders")

    if not isinstance(course_drive_root_value, str) or not course_drive_root_value.strip():
        raise ValueError("course_config.local.yaml is missing course_drive_root.")
    if not isinstance(folders, dict):
        raise ValueError("course_config.yaml is missing the folders mapping.")

    folder_name = folders.get(folder_key)
    if not isinstance(folder_name, str) or not folder_name.strip():
        raise ValueError(f"course_config.yaml is missing folders.{folder_key}.")

    return Path(course_drive_root_value).expanduser() / folder_name


def read_csv(path: Path) -> list[dict[str, str]]:
    with path.open("r", encoding="utf-8-sig", newline="") as handle:
        return list(csv.DictReader(handle))


def write_csv(path: Path, fieldnames: list[str], rows: list[dict[str, str]]) -> None:
    with path.open("w", encoding="utf-8-sig", newline="") as handle:
        writer = csv.DictWriter(handle, fieldnames=fieldnames)
        writer.writeheader()
        writer.writerows(rows)


def segment_files(root: Path) -> list[Path]:
    files = sorted(
        (root / VIDEO_SEGMENTS_DIR).glob("*_video_segments.csv"),
        key=lambda path: int(path.name.split("_", 1)[0]),
    )
    if not files:
        raise FileNotFoundError(f"No video segment CSV files found in {root / VIDEO_SEGMENTS_DIR}")
    return files


def parse_segment_row(path: Path, row_index: int, row: dict[str, str]) -> Segment:
    return Segment(
        source_csv=path,
        row_index=row_index,
        lecture_number=int(row["lecture_number"]),
        lecture_title=row["lecture_title"],
        video_number=int(row["video_number"]),
        video_title=row["video_title"],
        slide_start=int(row["original_slide_start"]),
        slide_end=int(row["original_slide_end"]),
        estimated_minutes=row["estimated_minutes"],
        main_concept_group=row["main_concept_group"],
        learning_objectives=row["learning_objectives"],
        logical_segment_reason=row["why_this_is_a_logical_segment"],
        stopping_point_reason=row["why_the_ending_slide_is_a_good_stopping_point"],
        recap=row["suggested_recap_slide"],
        quiz=row["suggested_quiz_question"],
        review_status=row["review_status"],
    )


def load_segment_rows(root: Path) -> dict[Path, tuple[list[str], list[dict[str, str]]]]:
    loaded: dict[Path, tuple[list[str], list[dict[str, str]]]] = {}
    for path in segment_files(root):
        with path.open("r", encoding="utf-8-sig", newline="") as handle:
            reader = csv.DictReader(handle)
            rows = list(reader)
            fieldnames = list(reader.fieldnames or [])
        if not fieldnames:
            raise ValueError(f"{path} has no CSV header.")
        if OUTPUT_FIELD not in fieldnames:
            fieldnames.append(OUTPUT_FIELD)
            for row in rows:
                row[OUTPUT_FIELD] = ""
        loaded[path] = (fieldnames, rows)
    return loaded


def approved_segments(
    loaded_rows: dict[Path, tuple[list[str], list[dict[str, str]]]]
) -> list[Segment]:
    segments: list[Segment] = []
    for path, (_fieldnames, rows) in loaded_rows.items():
        for row_index, row in enumerate(rows):
            if row.get("review_status") == APPROVED_STATUS:
                segments.append(parse_segment_row(path, row_index, row))
    return sorted(segments, key=lambda item: (item.lecture_number, item.video_number))


def csv_row_number(row_index: int) -> int:
    return row_index + 2


def row_matches_lecture_filter(row: dict[str, str], lecture_number: int | None) -> bool:
    if lecture_number is None:
        return True
    return row.get("lecture_number") == str(lecture_number)


def load_teaching_maps(root: Path) -> dict[int, dict[int, dict[str, str]]]:
    maps: dict[int, dict[int, dict[str, str]]] = {}
    for path in sorted((root / TEACHING_MAPS_DIR).glob("*_teaching_map.csv")):
        rows = read_csv(path)
        if not rows:
            continue
        lecture_number = int(rows[0]["lecture_number"])
        maps[lecture_number] = {int(row["slide_number"]): row for row in rows}
    return maps


def split_items(text: str, max_items: int | None = None) -> list[str]:
    parts = [
        part.strip(" .")
        for part in re.split(r";|\n", text)
        if part.strip(" .")
    ]
    if len(parts) <= 1:
        parts = [
            part.strip(" .")
            for part in re.split(r", and |, ", text)
            if part.strip(" .")
        ]
    if max_items is not None:
        parts = parts[:max_items]
    return parts or [text.strip()]


def shorten_text(text: str, max_length: int = 155) -> str:
    normalized = re.sub(r"\s+", " ", text).strip(" .")
    if len(normalized) <= max_length:
        return normalized
    shortened = normalized[: max_length - 3].rsplit(" ", 1)[0].strip(" .")
    return f"{shortened}..."


def recap_label(label: str, group: str) -> str:
    normalized_label = re.sub(r"\s+", " ", label).strip(" .")
    normalized_group = re.sub(r"\s+", " ", group).strip(" .")
    if normalized_label.lower() == normalized_group.lower():
        return ""
    prefix = f"{normalized_group}:"
    if normalized_label.lower().startswith(prefix.lower()):
        return normalized_label[len(prefix) :].strip(" .")
    return normalized_label


def segment_teaching_rows(
    segment: Segment,
    teaching_map: dict[int, dict[str, str]],
) -> list[dict[str, str]]:
    return [
        teaching_map[slide_number]
        for slide_number in range(segment.slide_start, segment.slide_end + 1)
        if slide_number in teaching_map
    ]


def recap_bullets(
    segment: Segment,
    teaching_map: dict[int, dict[str, str]],
) -> list[str]:
    rows = segment_teaching_rows(segment, teaching_map)
    bullets: list[str] = []

    grouped_rows: dict[str, list[dict[str, str]]] = {}
    for row in rows:
        group = row.get("concept_group", "").strip() or row.get("core_concept", "").strip()
        if not group:
            continue
        grouped_rows.setdefault(group, []).append(row)

    for group, group_rows in grouped_rows.items():
        labels: list[str] = []
        for row in group_rows:
            raw_label = row.get("slide_title", "").strip() or row.get("core_concept", "").strip()
            label = recap_label(raw_label, group)
            if label and label not in labels:
                labels.append(label)
        if not labels:
            bullets.append(shorten_text(group))
        elif len(labels) == 1:
            bullets.append(shorten_text(f"{group}: {labels[0]}"))
        else:
            bullets.append(shorten_text(f"{group}: {labels[0]} -> {labels[-1]}"))

    fallback_items = split_items(segment.recap, max_items=6)
    if len(bullets) < 3:
        for item in fallback_items:
            if len(bullets) >= 6:
                break
            if item and all(item.lower() not in bullet.lower() for bullet in bullets):
                bullets.append(shorten_text(item))

    return bullets[:6] or fallback_items


def slugify(value: str, max_length: int = 70) -> str:
    slug = re.sub(r"[^A-Za-z0-9]+", "_", value).strip("_").lower()
    slug = re.sub(r"_+", "_", slug)
    return (slug[:max_length].rstrip("_") or "segment")


def output_filename(segment: Segment) -> str:
    title_slug = slugify(segment.video_title)
    return f"lecture_{segment.lecture_number:02d}_video_{segment.video_number:02d}_{title_slug}.pptx"


def output_path_for(revised_root: Path, segment: Segment) -> Path:
    return revised_root / str(segment.lecture_number) / output_filename(segment)


def relative_output_name(revised_root: Path, output_path: Path) -> str:
    return output_path.relative_to(revised_root).as_posix()


def normalize_output_filter(value: str) -> str:
    return value.strip().replace("\\", "/")


def assert_safe_output(source_path: Path, revised_root: Path, output_path: Path) -> None:
    source_resolved = source_path.resolve(strict=False)
    output_resolved = output_path.resolve(strict=False)
    revised_resolved = revised_root.resolve(strict=False)

    if output_resolved == source_resolved:
        raise RuntimeError(f"Refusing to overwrite original PowerPoint: {source_path}")
    try:
        output_resolved.relative_to(revised_resolved)
    except ValueError as exc:
        raise RuntimeError(f"Output path is not inside revised slides folder: {output_path}") from exc


def delete_unwanted_slides(
    presentation: object,
    start: int,
    end: int,
    original_slide_count: int,
) -> None:
    keep = set(range(start, end + 1))
    slide_id_list = presentation.slides._sldIdLst  # python-pptx internal API.
    slide_ids = list(slide_id_list)

    for zero_based_index in reversed(range(len(slide_ids))):
        slide_number = zero_based_index + 1
        if slide_number in keep or slide_number > original_slide_count:
            continue
        slide_id = slide_ids[zero_based_index]
        presentation.part.drop_rel(slide_id.rId)
        slide_id_list.remove(slide_id)


def move_slide(presentation: object, old_index: int, new_index: int) -> None:
    slide_id_list = presentation.slides._sldIdLst  # python-pptx internal API.
    slide_ids = list(slide_id_list)
    slide_id = slide_ids[old_index]
    slide_id_list.remove(slide_id)
    slide_id_list.insert(new_index, slide_id)


def find_layout(presentation: object, preferred_names: list[str], fallback_index: int) -> object:
    for preferred_name in preferred_names:
        normalized = preferred_name.lower()
        for layout in presentation.slide_layouts:
            if layout.name.lower() == normalized:
                return layout
    if fallback_index < len(presentation.slide_layouts):
        return presentation.slide_layouts[fallback_index]
    return presentation.slide_layouts[0]


def clear_text_frame(text_frame: object) -> None:
    text_frame.clear()


def set_run_font_size(text_frame: object, size: int) -> None:
    if Pt is None:
        return
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(size)


def set_slide_title(
    slide: object,
    title: str,
    bounds: tuple[int, int, int, int] | None = None,
    font_size: int | None = None,
) -> None:
    if slide.shapes.title is not None:
        if bounds is not None:
            left, top, width, height = bounds
            slide.shapes.title.left = left
            slide.shapes.title.top = top
            slide.shapes.title.width = width
            slide.shapes.title.height = height
        slide.shapes.title.text = title
        if font_size is not None:
            set_run_font_size(slide.shapes.title.text_frame, font_size)
        return

    if Inches is None or Pt is None:
        raise RuntimeError("python-pptx is required to create title text boxes.")

    if bounds is None:
        bounds = (Inches(0.7), Inches(0.35), Inches(11.0), Inches(0.6))
    textbox = slide.shapes.add_textbox(*bounds)
    frame = textbox.text_frame
    frame.text = title
    set_run_font_size(frame, font_size or 28)


def body_shape(slide: object) -> object | None:
    if PP_PLACEHOLDER is not None:
        preferred_types = {
            PP_PLACEHOLDER.BODY,
            PP_PLACEHOLDER.OBJECT,
            PP_PLACEHOLDER.VERTICAL_BODY,
            PP_PLACEHOLDER.VERTICAL_OBJECT,
        }
        for placeholder in slide.placeholders:
            try:
                placeholder_type = placeholder.placeholder_format.type
            except ValueError:
                continue
            if placeholder_type in preferred_types and placeholder is not slide.shapes.title:
                return placeholder
    return None


def body_text_frame(
    slide: object,
    bounds: tuple[int, int, int, int] | None = None,
) -> object:
    shape = body_shape(slide)
    if shape is not None:
        if bounds is not None:
            left, top, width, height = bounds
            shape.left = left
            shape.top = top
            shape.width = width
            shape.height = height
        return shape.text_frame

    if Inches is None:
        raise RuntimeError("python-pptx is required to create body text boxes.")

    if bounds is None:
        bounds = (Inches(0.9), Inches(1.5), Inches(11.2), Inches(5.0))
    textbox = slide.shapes.add_textbox(*bounds)
    return textbox.text_frame


def set_bullets(
    slide: object,
    bullets: list[str],
    font_size: int = 22,
    bounds: tuple[int, int, int, int] | None = None,
) -> None:
    frame = body_text_frame(slide, bounds)
    clear_text_frame(frame)

    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        set_run_font_size(frame, font_size)


def generated_title_bounds(presentation: object) -> tuple[int, int, int, int]:
    if Inches is None:
        raise RuntimeError("python-pptx is required to place generated text.")
    return (
        Inches(0.75),
        Inches(0.35),
        presentation.slide_width - Inches(1.5),
        Inches(0.85),
    )


def generated_body_bounds(presentation: object) -> tuple[int, int, int, int]:
    if Inches is None:
        raise RuntimeError("python-pptx is required to place generated text.")
    top = Inches(1.45)
    return (
        Inches(0.9),
        top,
        presentation.slide_width - Inches(1.8),
        presentation.slide_height - top - Inches(0.7),
    )


def add_notes_body_placeholder(notes_slide: object) -> object:
    if parse_xml is None or nsdecls is None:
        raise RuntimeError("python-pptx XML helpers are required to create speaker notes.")

    sp_tree = notes_slide._element.xpath("./p:cSld/p:spTree")[0]
    shape_ids: list[int] = []
    for element in sp_tree.xpath(".//p:cNvPr"):
        shape_id = element.get("id")
        if shape_id and shape_id.isdigit():
            shape_ids.append(int(shape_id))
    next_shape_id = max(shape_ids, default=1) + 1

    notes_shape = parse_xml(
        f"""<p:sp {nsdecls("p", "a")}>
  <p:nvSpPr>
    <p:cNvPr id="{next_shape_id}" name="Generated Notes Placeholder"/>
    <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
    <p:nvPr><p:ph type="body" idx="3" sz="quarter"/></p:nvPr>
  </p:nvSpPr>
  <p:spPr/>
  <p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>
</p:sp>"""
    )
    sp_tree.append(notes_shape)
    return notes_slide.notes_text_frame


def notes_text_frame(slide: object) -> object:
    notes_slide = slide.notes_slide
    frame = notes_slide.notes_text_frame
    if frame is not None:
        return frame
    return add_notes_body_placeholder(notes_slide)


def set_notes(slide: object, notes: str, append: bool = False) -> None:
    frame = notes_text_frame(slide)
    existing = frame.text.strip()
    if append and existing:
        frame.text = f"{existing}\n\n{notes}"
    else:
        frame.text = notes


def add_title_slide(presentation: object, segment: Segment) -> object:
    layout = find_layout(
        presentation,
        ["Title and Content", "Title Only", "Section Header", "Title Slide"],
        1,
    )
    slide = presentation.slides.add_slide(layout)
    set_slide_title(
        slide,
        segment.video_title,
        bounds=generated_title_bounds(presentation),
        font_size=32,
    )
    set_bullets(
        slide,
        [
            f"Lecture {segment.lecture_number}: {segment.lecture_title}",
            f"Slides {segment.slide_start}-{segment.slide_end} | estimated teaching time: {segment.estimated_minutes} minutes",
            "By the end, you should be able to:",
            *split_items(segment.learning_objectives, max_items=4),
        ],
        font_size=20,
        bounds=generated_body_bounds(presentation),
    )
    set_notes(
        slide,
        "\n".join(
            [
                SCRIPT_MARKER,
                f"Open by framing this video as: {segment.video_title}.",
                f"Learning objectives: {segment.learning_objectives}",
                f"Segment rationale: {segment.logical_segment_reason}",
            ]
        ),
    )
    return slide


def add_recap_slide(
    presentation: object,
    segment: Segment,
    teaching_map: dict[int, dict[str, str]],
) -> object:
    layout = find_layout(presentation, ["Title and Content", "Title Only"], 1)
    slide = presentation.slides.add_slide(layout)
    set_slide_title(
        slide,
        "Recap",
        bounds=generated_title_bounds(presentation),
        font_size=32,
    )
    bullets = recap_bullets(segment, teaching_map)
    set_bullets(
        slide,
        bullets,
        font_size=20,
        bounds=generated_body_bounds(presentation),
    )
    set_notes(
        slide,
        "\n".join(
            [
                SCRIPT_MARKER,
                "Use this slide as a short verbal synthesis before the checkpoint.",
                f"Recap focus: {segment.recap}",
                "Expanded recap bullets:",
                *bullets,
                f"Stopping point: {segment.stopping_point_reason}",
            ]
        ),
    )
    return slide


def add_quiz_slide(presentation: object, segment: Segment) -> object:
    layout = find_layout(presentation, ["Title and Content", "Title Only"], 1)
    slide = presentation.slides.add_slide(layout)
    set_slide_title(
        slide,
        "Checkpoint",
        bounds=generated_title_bounds(presentation),
        font_size=32,
    )
    set_bullets(
        slide,
        [segment.quiz],
        font_size=22,
        bounds=generated_body_bounds(presentation),
    )
    set_notes(
        slide,
        "\n".join(
            [
                SCRIPT_MARKER,
                "Pause here and ask students to answer before continuing.",
                f"Checkpoint question: {segment.quiz}",
                "Use the response to confirm whether the recap concepts are clear.",
            ]
        ),
    )
    return slide


def outline_for_source_slide(
    slide_number: int,
    teaching_map: dict[int, dict[str, str]],
) -> str:
    row = teaching_map.get(slide_number)
    if row is None:
        return "\n".join(
            [
                SCRIPT_MARKER,
                f"Original slide {slide_number}.",
                "No teaching-map row was found for this slide.",
            ]
        )

    lines = [
        SCRIPT_MARKER,
        f"Slide title: {row['slide_title']}",
        f"Core concept: {row['core_concept']}",
        f"Role in lesson: {row['slide_role']}",
        f"Concept group: {row['concept_group']}",
        f"Teaching note: {row['notes_for_future_script']}",
    ]
    if row.get("prerequisite_concept"):
        lines.append(f"Connect from: {row['prerequisite_concept']}")
    if row.get("next_concept"):
        lines.append(f"Lead into: {row['next_concept']}")
    return "\n".join(lines)


def add_source_slide_notes(
    presentation: object,
    segment: Segment,
    teaching_maps: dict[int, dict[int, dict[str, str]]],
) -> None:
    teaching_map = teaching_maps.get(segment.lecture_number, {})
    for slide_number in range(segment.slide_start, segment.slide_end + 1):
        slide = presentation.slides[slide_number - 1]
        set_notes(slide, outline_for_source_slide(slide_number, teaching_map), append=True)


def build_segment_deck(
    source_path: Path,
    output_path: Path,
    segment: Segment,
    teaching_maps: dict[int, dict[int, dict[str, str]]],
) -> None:
    with open_presentation_readonly(source_path) as (presentation, _repair_info):
        if segment.slide_end > len(presentation.slides):
            raise ValueError(
                f"{source_path.name} has {len(presentation.slides)} slides, "
                f"but segment {segment.video_number} ends at {segment.slide_end}."
            )

        original_slide_count = len(presentation.slides)
        teaching_map = teaching_maps.get(segment.lecture_number, {})
        add_title_slide(presentation, segment)
        add_recap_slide(presentation, segment, teaching_map)
        add_quiz_slide(presentation, segment)
        add_source_slide_notes(presentation, segment, teaching_maps)
        delete_unwanted_slides(
            presentation,
            segment.slide_start,
            segment.slide_end,
            original_slide_count,
        )
        move_slide(presentation, len(presentation.slides) - 3, 0)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(output_path)


def update_source_csv_rows(
    loaded_rows: dict[Path, tuple[list[str], list[dict[str, str]]]],
    segment: Segment,
    revised_root: Path,
    output_path: Path,
) -> None:
    _fieldnames, rows = loaded_rows[segment.source_csv]
    rows[segment.row_index][OUTPUT_FIELD] = relative_output_name(revised_root, output_path)


def write_updated_segment_csvs(
    loaded_rows: dict[Path, tuple[list[str], list[dict[str, str]]]]
) -> None:
    for path, (fieldnames, rows) in loaded_rows.items():
        write_csv(path, fieldnames, rows)


def deck_map(original_slides_folder: Path) -> dict[int, LectureDeck]:
    return {deck.lecture_number: deck for deck in discover_lecture_decks(original_slides_folder)}


def print_dry_run_report(
    loaded_rows: dict[Path, tuple[list[str], list[dict[str, str]]]],
    segments: list[Segment],
    decks: dict[int, LectureDeck],
    revised_slides_folder: Path,
    lecture_number_filter: int | None,
) -> None:
    print("Dry run report")
    print("Segment CSV files that would be read:")
    for path, (_fieldnames, rows) in sorted(
        loaded_rows.items(),
        key=lambda item: int(item[0].name.split("_", 1)[0]),
    ):
        print(f"- {path} ({len(rows)} rows)")

    print()
    print("Rows approved_for_pptx that would be processed:")
    for segment in segments:
        deck = decks.get(segment.lecture_number)
        if deck is None:
            raise FileNotFoundError(
                f"No original PowerPoint found for lecture {segment.lecture_number}."
            )
        output_path = output_path_for(revised_slides_folder, segment)
        assert_safe_output(deck.path, revised_slides_folder, output_path)
        print(
            f"- {segment.source_csv.name} row {csv_row_number(segment.row_index)}: "
            f"lecture {segment.lecture_number} video {segment.video_number} "
            f"({segment.video_title})"
        )

    print()
    print("Original PowerPoint files that would be read:")
    for source_path in sorted(
        {decks[segment.lecture_number].path for segment in segments},
        key=lambda path: path.name.lower(),
    ):
        print(f"- {source_path}")

    print()
    print("Output PowerPoint files that would be created:")
    for segment in segments:
        output_path = output_path_for(revised_slides_folder, segment)
        print(
            f"- lecture {segment.lecture_number} video {segment.video_number}: "
            f"{output_path}"
        )

    print()
    print("Rows that would be skipped:")
    skipped_count = 0
    for path, (_fieldnames, rows) in sorted(
        loaded_rows.items(),
        key=lambda item: int(item[0].name.split("_", 1)[0]),
    ):
        for row_index, row in enumerate(rows):
            status = row.get("review_status", "")
            if not row_matches_lecture_filter(row, lecture_number_filter):
                skipped_count += 1
                print(
                    f"- {path.name} row {csv_row_number(row_index)}: "
                    f"lecture {row.get('lecture_number', '')} "
                    f"video {row.get('video_number', '')} skipped by --lecture-number"
                )
            elif status != APPROVED_STATUS:
                skipped_count += 1
                print(
                    f"- {path.name} row {csv_row_number(row_index)}: "
                    f"lecture {row.get('lecture_number', '')} "
                    f"video {row.get('video_number', '')} "
                    f"status={status or '<missing>'}"
                )
    if skipped_count == 0:
        print("- None")


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Generate revised PowerPoint decks for approved video segments."
    )
    parser.add_argument(
        "--dry-run",
        action="store_true",
        help="Validate inputs and print planned outputs without writing PPTX files or CSVs.",
    )
    parser.add_argument(
        "--overwrite",
        action="store_true",
        help="Overwrite existing generated segment decks in 02_Revised_Slides.",
    )
    parser.add_argument(
        "--lecture-number",
        type=int,
        help="Process only one lecture number.",
    )
    parser.add_argument(
        "--output-pptx",
        action="append",
        default=[],
        help=(
            "Process only the segment whose recorded output_pptx matches this "
            "relative path. Can be passed more than once."
        ),
    )
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    root = project_root()

    try:
        if PP_PLACEHOLDER is None or Inches is None or Pt is None:
            raise RuntimeError(
                "python-pptx is required. Install dependencies with: pip install -r requirements.txt"
            )

        original_slides_folder = configured_folder(root, "original_slides")
        revised_slides_folder = configured_folder(root, "revised_slides")
        decks = deck_map(original_slides_folder)
        loaded_rows = load_segment_rows(root)
        teaching_maps = load_teaching_maps(root)
        segments = approved_segments(loaded_rows)
        if args.lecture_number is not None:
            segments = [
                segment
                for segment in segments
                if segment.lecture_number == args.lecture_number
            ]
        if args.output_pptx:
            requested_outputs = {
                normalize_output_filter(value) for value in args.output_pptx
            }
            segments = [
                segment
                for segment in segments
                if relative_output_name(
                    revised_slides_folder,
                    output_path_for(revised_slides_folder, segment),
                )
                in requested_outputs
            ]
        if not segments:
            raise ValueError("No approved video segments found to generate.")

        if args.dry_run:
            print_dry_run_report(
                loaded_rows,
                segments,
                decks,
                revised_slides_folder,
                args.lecture_number,
            )
            print()
            print("Segment deck generation complete")
            print(f"Original slides folder: {original_slides_folder}")
            print(f"Revised slides folder:  {revised_slides_folder}")
            print(f"Approved segments:      {len(segments)}")
            print("Dry run:                no PPTX files or CSV files were written.")
            print("PowerPoint safety: original files were discovered only; no decks were opened.")
            return 0

        created = 0
        skipped_existing = 0
        for segment in segments:
            deck = decks.get(segment.lecture_number)
            if deck is None:
                raise FileNotFoundError(
                    f"No original PowerPoint found for lecture {segment.lecture_number}."
                )

            output_path = output_path_for(revised_slides_folder, segment)
            assert_safe_output(deck.path, revised_slides_folder, output_path)

            if output_path.exists() and not args.overwrite:
                skipped_existing += 1
                update_source_csv_rows(loaded_rows, segment, revised_slides_folder, output_path)
                continue

            build_segment_deck(deck.path, output_path, segment, teaching_maps)
            update_source_csv_rows(loaded_rows, segment, revised_slides_folder, output_path)
            created += 1

        write_updated_segment_csvs(loaded_rows)

    except (OSError, RuntimeError, ValueError, KeyError, YAMLError, PptxReadError) as exc:
        print("Segment deck generation failed.")
        print(f"Error: {exc}")
        return 1

    print("Segment deck generation complete")
    print(f"Original slides folder: {original_slides_folder}")
    print(f"Revised slides folder:  {revised_slides_folder}")
    print(f"Approved segments:      {len(segments)}")
    print(f"Decks created:          {created}")
    print(f"Existing decks skipped: {skipped_existing}")
    print(f"CSV field updated:      {OUTPUT_FIELD}")
    print("PowerPoint safety: original files were opened read-only and never overwritten.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
